#!/usr/bin/env python3
"""
Gamma API调用模块
用于根据用户提示生成PPTX文件
"""

import os
import requests
import json
import time
import platform
from pathlib import Path
from typing import Optional, Dict, Any
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches

# 加载环境变量
load_dotenv()


class GammaAPI:
    """Gamma API客户端类"""

    def __init__(self, api_key: Optional[str] = None):
        """
        初始化Gamma API客户端

        Args:
            api_key: Gamma API密钥，如果不提供则从环境变量读取
        """
        self.api_key = api_key or os.getenv("GAMMA_API_KEY") or os.getenv("gamma_api_key")
        self.base_url = "https://public-api.gamma.app/v0.2"

        if not self.api_key:
            raise ValueError("未找到GAMMA_API_KEY或gamma_api_key，请确保在.env文件中设置了该环境变量")

    def _read_prompt_file(self, prompt_file: str = "input/prompt.txt") -> str:
        """
        读取提示词文件内容

        Args:
            prompt_file: 提示词文件路径

        Returns:
            提示词内容

        Raises:
            FileNotFoundError: 文件不存在
        """
        file_path = Path(prompt_file)
        if not file_path.exists():
            raise FileNotFoundError(f"提示词文件不存在: {file_path}")

        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read().strip()

            if not content:
                raise ValueError(f"提示词文件为空: {file_path}")

            return content
        except Exception as e:
            raise Exception(f"读取提示词文件时发生错误: {str(e)}")

    def _add_notes_to_pptx(self, pptx_path: str, prompt_file: str = "input/prompt.txt") -> bool:
        """
        为下载的PPTX文件添加备注

        Args:
            pptx_path: PPTX文件路径
            prompt_file: 提示词文件路径

        Returns:
            是否成功添加备注
        """
        try:
            # 读取并解析prompt.txt文件
            prompt_content = self._read_prompt_file(prompt_file)
            # 使用"\n---\n"分隔每一页的内容
            pages_content = prompt_content.split('\n---\n')

            # 清理每页内容（去掉首尾空白）
            pages_content = [page.strip() for page in pages_content if page.strip()]

            print(f"📄 从 prompt.txt 解析出 {len(pages_content)} 页内容")

            # 打开PPTX文件
            prs = Presentation(pptx_path)
            total_slides = len(prs.slides)
            print(f"📑 PPTX文件共有 {total_slides} 页")

            # 为每一页幻灯片添加备注
            notes_added = 0
            for i, slide in enumerate(prs.slides):
                if i < len(pages_content):
                    notes_content = pages_content[i]

                    # 获取备注区域
                    if slide.has_notes_slide:
                        notes_slide = slide.notes_slide
                        # 清除现有备注内容
                        for shape in notes_slide.shapes:
                            if shape.has_text_frame:
                                shape.text_frame.clear()

                        # 添加新的备注内容
                        notes_placeholder = None
                        for shape in notes_slide.shapes:
                            if shape.has_text_frame and hasattr(shape, 'placeholder_format'):
                                # 寻找备注占位符或任何可以用于备注的文本框
                                if shape.placeholder_format.type == 2:  # 2 = NOTES
                                    notes_placeholder = shape
                                    break
                                elif notes_placeholder is None and shape.name and 'Notes' in shape.name:
                                    notes_placeholder = shape

                        # 如果找到备注占位符，就使用它；否则使用第一个文本框
                        target_shape = notes_placeholder
                        if target_shape is None:
                            for shape in notes_slide.shapes:
                                if shape.has_text_frame:
                                    target_shape = shape
                                    break

                        if target_shape:
                            text_frame = target_shape.text_frame
                            text_frame.text = notes_content
                    else:
                        # 如果没有备注页，创建一个
                        notes_slide = slide.notes_slide

                        # 添加新的备注内容（逻辑同上）
                        notes_placeholder = None
                        for shape in notes_slide.shapes:
                            if shape.has_text_frame and hasattr(shape, 'placeholder_format'):
                                # 寻找备注占位符或任何可以用于备注的文本框
                                if shape.placeholder_format.type == 2:  # 2 = NOTES
                                    notes_placeholder = shape
                                    break
                                elif notes_placeholder is None and shape.name and 'Notes' in shape.name:
                                    notes_placeholder = shape

                        # 如果找到备注占位符，就使用它；否则使用第一个文本框
                        target_shape = notes_placeholder
                        if target_shape is None:
                            for shape in notes_slide.shapes:
                                if shape.has_text_frame:
                                    target_shape = shape
                                    break

                        if target_shape:
                            text_frame = target_shape.text_frame
                            text_frame.text = notes_content

                    print(f"✅ 已为第 {i+1} 页添加备注（内容长度: {len(notes_content)} 字符）")
                    notes_added += 1
                else:
                    print(f"⚠️ 第 {i+1} 页无对应备注内容（prompt.txt 只有 {len(pages_content)} 页）")

            # 保存修改后的PPTX文件
            prs.save(pptx_path)
            print(f"✅ 成功为PPTX文件添加备注: {pptx_path}")
            print(f"📊 总计: PPTX {total_slides} 页，添加备注 {notes_added} 页")

            return True

        except Exception as e:
            print(f"❌ 添加备注时发生错误: {str(e)}")
            return False

    def _ppt_to_pdf(self, input_path: str, output_path: str) -> bool:
        """
        将PPTX文件转换为PDF文件（仅支持Windows系统）

        Args:
            input_path: PPTX文件路径
            output_path: PDF输出路径

        Returns:
            转换是否成功
        """
        try:
            # 仅在Windows系统下运行
            if platform.system() != "Windows":
                print("⚠️ PPT转PDF功能仅在Windows系统下可用")
                return False

            import win32com.client
            
            print(f"🔄 开始转换PPTX为PDF: {input_path} -> {output_path}")
            
            # 启动WPS应用
            wps = win32com.client.Dispatch("Kwpp.Application")
            
            try:
                # 打开演示文稿
                presentation = wps.Presentations.Open(input_path, WithWindow=False)
                
                # 保存为PDF格式（32表示PDF格式）
                presentation.SaveAs(output_path, 32)
                
                # 关闭演示文稿
                presentation.Close()
                
                print(f"✅ PPTX已成功转换为PDF: {output_path}")
                return True
                
            finally:
                # 确保WPS应用退出
                wps.Quit()
                
        except ImportError:
            print("❌ 缺少win32com模块，请安装pywin32：pip install pywin32")
            return False
        except Exception as e:
            print(f"❌ PPT转PDF过程中发生错误: {str(e)}")
            return False

    def _prepare_generation_payload(self, input_text: str, **kwargs) -> Dict[str, Any]:
        """
        准备生成PPT的请求载荷

        Args:
            input_text: 输入文本内容
            **kwargs: 其他可选参数

        Returns:
            请求载荷字典
        """
        # 默认配置
        default_config = {
            "inputText": input_text,
            "textMode": "generate",
            "format": "presentation",
            "themeName": "企业汇报模版",
            "numCards": 10,
            "cardSplit": "auto",
            "exportAs": "pptx",
            "textOptions": {
                "amount": "brief",
                "tone": "professional",
                "audience": "college students",
                "language": "zh-cn"
            },
            "imageOptions": {
                "source": "aiGenerated",
                "model": "",
                "style": "插画风格"
            },
            "cardOptions": {"dimensions": "16x9"},
            "sharingOptions": {
                "workspaceAccess": "fullAccess",
                "externalAccess": "noAccess"
            }
        }

        # 合并用户提供的参数
        payload = {**default_config, **kwargs}

        # 确保必填字段存在
        if "inputText" not in payload or not payload["inputText"]:
            raise ValueError("inputText是必填字段")

        return payload

    def _download_file(self, download_url: str, output_path: str, prompt_file: str = "input/prompt.txt") -> bool:
        """
        下载生成的文件并为PPTX添加备注

        Args:
            download_url: 下载链接
            output_path: 输出文件路径
            prompt_file: 提示词文件路径

        Returns:
            下载是否成功
        """
        try:
            response = requests.get(download_url, timeout=300)
            # 确保输出目录存在
            output_file = Path(output_path)
            output_file.parent.mkdir(parents=True, exist_ok=True)

            # 写入文件
            with open(output_file, 'wb') as f:
                f.write(response.content)

            # 如果是PPTX文件，添加备注
            if output_path.lower().endswith('.pptx'):
                print("📝 开始为PPTX文件添加备注...")
                success = self._add_notes_to_pptx(output_path, prompt_file)
                if success:
                    print("✅ 备注添加完成")
                else:
                    print("⚠️ 备注添加失败，但文件下载成功")

                # 如果在Windows系统下，将PPTX转换为PDF
                if platform.system() == "Windows":
                    print("🔄 检测到Windows系统，开始转换PPTX为PDF...")
                    try:
                        # 获取绝对路径
                        pptx_abs_path = os.path.abspath(output_path)
                        
                        # 生成PDF文件路径（相同目录，扩展名改为.pdf）
                        pdf_path = output_path.replace('.pptx', '.pdf')
                        pdf_abs_path = os.path.abspath(pdf_path)
                        
                        # 确保PDF输出目录存在
                        pdf_dir = os.path.dirname(pdf_abs_path)
                        os.makedirs(pdf_dir, exist_ok=True)
                        
                        # 使用内置方法转换PPTX为PDF
                        success = self._ppt_to_pdf(pptx_abs_path, pdf_abs_path)
                        if not success:
                            print("⚠️ PDF转换失败，但PPTX文件下载成功")
                        
                    except Exception as pdf_error:
                        print(f"⚠️ PDF转换过程中发生错误: {str(pdf_error)}")

            return True
        except Exception as e:
            print(f"下载文件时发生错误: {str(e)}")
            if hasattr(response, 'status_code'):
                print(f"HTTP状态码: {response.status_code}")
                print(f"响应头: {response.headers}")
            return False

    def generate_pptx(
        self,
        prompt_file: str = "input/prompt.txt",
        output_dir: str = "input",
        **kwargs
    ) -> Optional[str]:
        """
        根据提示词文件生成PPTX文件

        Args:
            prompt_file: 提示词文件路径
            output_dir: 输出目录（默认为input目录）
            **kwargs: 其他可选参数

        Returns:
            生成的PPTX文件路径，如果失败则返回None
        """
        try:
            # 读取提示词
            print("📖 读取提示词文件...")
            input_text = self._read_prompt_file(prompt_file)
            print(f"✅ 成功读取提示词，长度: {len(input_text)} 字符")

            # 准备请求载荷
            print("🔧 准备API请求参数...")
            payload = self._prepare_generation_payload(input_text, **kwargs)

            # 准备请求头
            headers = {
                "accept": "application/json",
                "Content-Type": "application/json",
                "X-API-KEY": self.api_key
            }

            # 发送生成请求
            url = f"{self.base_url}/generations"
            print("🚀 发送PPT生成请求...")

            response = requests.post(url, json=payload, headers=headers, timeout=60)

            # 打印详细的错误信息
            if not response.ok:
                print(f"❌ API请求失败，状态码: {response.status_code}")
                print(f"❌ 响应内容: {response.text}")
                print(f"❌ 请求数据: {payload}")

            response.raise_for_status()

            result = response.json()
            print("✅ PPT生成任务已提交")

            # 获取generationId
            generation_id = result.get("generationId")
            if not generation_id:
                raise ValueError("API响应中未找到generationId")

            print(f"📋 生成ID: {generation_id}")

            # 轮询任务状态 - 每30秒检查一次
            max_attempts = 60  # 最多等待60次，每次30秒 = 30分钟
            attempt = 0
            gamma_url = None

            while attempt < max_attempts:
                print(f"⏳ 检查任务状态... (尝试 {attempt + 1}/{max_attempts})")

                # 调用状态检查接口 - 根据文档应该是GET请求
                status_url = f"{self.base_url}/generations/{generation_id}"
                status_response = requests.get(status_url, headers=headers, timeout=30)

                # 打印详细的错误信息
                if not status_response.ok:
                    print(f"❌ 状态检查请求失败，状态码: {status_response.status_code}")
                    print(f"❌ 响应内容: {status_response.text}")
                    print(f"❌ 请求URL: {status_url}")

                status_response.raise_for_status()

                # 直接获取状态数据
                matching_generation = status_response.json()

                if not matching_generation:
                    print(f"⚠️ 未找到generationId为 {generation_id} 的任务")
                    time.sleep(30)
                    attempt += 1
                    continue

                status = matching_generation.get("status")
                print(f"🔄 任务状态: {status}")

                if status == "completed":
                    print("✅ PPT生成完成")
                    gamma_url = matching_generation.get("exportUrl")
                    credits_info = matching_generation.get("credits", {})
                    print(f"💰 消耗积分: {credits_info.get('deducted', 'N/A')}, 剩余积分: {credits_info.get('remaining', 'N/A')}")
                    break
                elif status == "failed":
                    error_msg = matching_generation.get("error", "未知错误")
                    raise Exception(f"PPT生成失败: {error_msg}")
                elif status in ["processing", "pending", "queued"]:
                    print(f"🔄 任务进行中... 状态: {status}")
                    time.sleep(30)  # 等待30秒
                else:
                    print(f"⚠️ 未知状态: {status}")
                    time.sleep(30)

                attempt += 1

            if attempt >= max_attempts:
                raise Exception("PPT生成超时，请稍后手动检查")

            if not gamma_url:
                raise Exception("生成完成但未找到gammaUrl")

            print(f"🔗 Gamma链接: {gamma_url}")

            # 使用导出API下载PPTX文件
            export_url = f"{gamma_url}"

            # 确保输出目录存在
            output_path = Path(output_dir)
            output_path.mkdir(parents=True, exist_ok=True)

            # 生成输出文件名
            timestamp = int(time.time())
            output_filename = f"generated_ppt_{timestamp}.pptx"
            full_output_path = output_path / output_filename

            print("📥 从Gamma导出PPTX文件...")
            success = self._download_file(export_url, str(full_output_path), prompt_file)

            if success:
                print(f"✅ PPTX文件已保存到: {full_output_path}")
                return str(full_output_path)
            else:
                print("❌ 文件下载失败")
                return None

        except Exception as e:
            print(f"❌ PPT生成过程中发生错误: {str(e)}")
            return None

    def get_theme_list(self) -> Optional[Dict[str, Any]]:
        """
        获取可用的主题列表

        Returns:
            主题列表数据
        """
        try:
            headers = {
                "accept": "application/json",
                "X-API-KEY": self.api_key
            }

            url = f"{self.base_url}/themes"
            response = requests.get(url, headers=headers, timeout=30)
            response.raise_for_status()

            return response.json()

        except Exception as e:
            print(f"❌ 获取主题列表时发生错误: {str(e)}")
            return None


def generate_pptx_from_prompt(
    prompt_file: str = "input/prompt.txt",
    output_dir: str = "output",
    **kwargs
) -> Optional[str]:
    """
    便捷函数：根据提示词文件生成PPTX

    Args:
        prompt_file: 提示词文件路径
        output_dir: 输出目录
        **kwargs: 其他可选参数

    Returns:
        生成的PPTX文件路径，如果失败则返回None
    """
    try:
        api = GammaAPI()
        return api.generate_pptx(prompt_file, output_dir, **kwargs)
    except Exception as e:
        print(f"❌ 生成PPTX时发生错误: {str(e)}")
        return None


if __name__ == "__main__":
    # 测试用例
    print("🧪 测试Gamma API功能...")

    # 确保提示词文件存在
    prompt_file = "input/prompt.txt"
    if not Path(prompt_file).exists():
        print(f"❌ 测试失败：提示词文件不存在 {prompt_file}")
        print("请创建input/prompt.txt文件并添加您要生成PPT的提示词内容")
        exit(1)

    # 测试生成PPTX
    result = generate_pptx_from_prompt(
        prompt_file=prompt_file,
        output_dir="output",
        themeName="企业汇报模版",
        numCards=8,
        additionalInstructions="创建一个专业的演示文稿，包含清晰的标题和结构化的内容"
    )

    if result:
        print(f"🎉 测试成功！生成的PPTX文件: {result}")
    else:
        print("❌ 测试失败：无法生成PPTX文件")