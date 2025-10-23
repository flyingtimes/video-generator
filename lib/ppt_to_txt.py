"""
PPT/PPTX文件备注提取工具
从PowerPoint文件中提取每一页的备注内容，并保存为txt文件
"""

import os
import sys
from pathlib import Path
from typing import List, Optional

# 添加项目根目录到Python路径，以便导入其他模块
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

try:
    from pptx import Presentation
except ImportError:
    print("错误: python-pptx库未安装，请运行: uv add python-pptx")
    sys.exit(1)


def extract_notes_from_pptx(pptx_path: str, output_dir: str) -> bool:
    """
    从PPTX文件中提取备注并保存到txt文件

    Args:
        pptx_path: PPTX文件路径
        output_dir: 输出目录路径

    Returns:
        bool: 是否成功提取
    """
    try:
        # 确保输出目录存在
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)

        # 加载演示文稿
        prs = Presentation(pptx_path)

        print(f"正在处理文件: {pptx_path}")
        print(f"总页数: {len(prs.slides)}")

        extracted_count = 0

        for i, slide in enumerate(prs.slides, 1):
            # 获取幻灯片的备注
            notes_text = ""

            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text += notes_slide.notes_text_frame.text.strip()

                # 去除空行并整理格式
                if notes_text:
                    lines = notes_text.split('\n')
                    # 过滤掉空行
                    notes_text  = '\n'.join([line.strip() for line in lines if line.strip()])


            # 生成输出文件路径
            output_file = output_path / f"{i}.txt"

            # 写入文件（如果无备注则写入空文件）
            with open(output_file, 'w', encoding='utf-8') as f:
                if notes_text:
                    f.write(notes_text)
                    print(f"  页 {i}: 已提取备注 ({len(notes_text)} 字符)")
                    extracted_count += 1
                else:
                    f.write("")  # 创建空文件
                    print(f"  页 {i}: 无备注内容")

        print(f"提取完成! 共处理 {len(prs.slides)} 页，其中 {extracted_count} 页包含备注")
        print(f"输出文件已保存到: {output_dir}")

        return True

    except Exception as e:
        print(f"错误: 处理文件 {pptx_path} 时发生异常: {str(e)}")
        return False


def find_pptx_file(input_dir: str, base_name: str) -> Optional[str]:
    """
    在输入目录中查找PPTX/PPT文件

    Args:
        input_dir: 输入目录
        base_name: 基础文件名（不含扩展名）

    Returns:
        Optional[str]: 找到的PPTX/PPT文件路径，如果没找到则返回None
    """
    input_path = Path(input_dir)

    # 按优先级查找文件
    extensions = ['.pptx', '.ppt']

    for ext in extensions:
        pptx_file = input_path / (base_name + ext)
        if pptx_file.exists():
            return str(pptx_file)

    return None


def main():
    """主函数"""
    if len(sys.argv) < 3:
        print("用法: python ppt_to_txt.py <pptx文件路径> <输出目录>")
        print("示例: python ppt_to_txt.py input/xh.pptx slides")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = sys.argv[2]

    if not os.path.exists(pptx_path):
        print(f"错误: 文件 {pptx_path} 不存在")
        sys.exit(1)

    success = extract_notes_from_pptx(pptx_path, output_dir)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()